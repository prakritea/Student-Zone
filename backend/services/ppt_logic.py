# backend/services/ppt_logic.py

import os
from typing import List
from pptx import Presentation
from pptx.util import Inches, Pt

# If using OpenAI for content generation
import openai
openai.api_key = os.getenv("OPENAI_API_KEY")

from models.schemas import PPTRequest, SlideContent

def generate_slide_content(topic: str, num_slides: int) -> List[SlideContent]:
    """
    Use LLM to generate titles + bullets for each slide given a topic.
    """
    prompt = f"""
You are an assistant that creates PowerPoint slide outlines.
Generate {num_slides} slides for the topic: "{topic}".
For each slide, give a title and 3‑5 bullet points.
Return in JSON format, like:
[
  {{
    "title": "Slide 1 title",
    "bullets": ["point 1", "point 2", "point 3"]
  }},
  ...
]
"""
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that helps generate slide content."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7,
        max_tokens=500
    )
    content = response.choices[0].message.content.strip()
    import json
    try:
        slides = json.loads(content)
    except json.JSONDecodeError:
        # fallback: simple parsing
        slides = []
        # Simple fallback: split by lines etc. (not shown here)
    # Each slide: ensure title and bullets
    results: List[SlideContent] = []
    for s in slides:
        title = s.get("title", "").strip()
        bullets = s.get("bullets", [])
        if title and isinstance(bullets, list) and bullets:
            results.append(SlideContent(title=title, bullets=[b.strip() for b in bullets if isinstance(b, str) and b.strip()]))
    return results

def build_ppt_from_content(slides: List[SlideContent], template_id: str, output_path: str) -> str:
    """
    Using a chosen template, build a pptx file with content and save to output_path.
    Returns the path of generated ppt file.
    """
    template_path = os.path.join("templates", f"{template_id}.pptx")
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template {template_id} not found.")
    
    prs = Presentation(template_path)
    
    # Choose which slide layout index to use for slides; depends on template
    # For example, assume template layout 1 is Title + Content
    slide_layout = prs.slide_layouts[1]  # adjust as per template

    for slide_content in slides:
        slide = prs.slides.add_slide(slide_layout)
        # set title
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_content.title

        # find content placeholder (usually the first placeholder that's not title)
        # Then add bullets
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            # ignore title
            if shp == title_placeholder:
                continue
            tf = shp.text_frame
            tf.clear()
            for b in slide_content.bullets:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0
            break

    # Save
    prs.save(output_path)
    return output_path
